from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
ASSETS_DIR = PRESENTATION_DIR / "assets"
OUTPUT_PATH = PRESENTATION_DIR / "Plastic_Degradation_Project_Presentation.pptx"

BG = "F6F7FB"
WHITE = "FFFFFF"
INK = "162033"
MUTED = "5F6B7A"
ACCENT = "3257FF"
ACCENT_SOFT = "E8EDFF"
TEAL = "0F766E"
ORANGE = "E76F51"
BORDER = "D6DCE8"


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def set_bg(slide, color: str = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_box(slide, left: float, top: float, width: float, height: float, *, fill_color: str, line_color: str | None = None, radius: bool = True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(1)
    return shape


def add_text(slide, left: float, top: float, width: float, height: float, text: str, *, size: float = 18, color: str = INK, bold: bool = False, align=PP_ALIGN.LEFT, name: str | None = None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, 0.7, 0.35, 10.5, 0.5, title, size=26, bold=True)
    add_box(slide, 0.72, 0.97, 1.7, 0.06, fill_color=ACCENT, radius=False)
    if subtitle:
        add_text(slide, 0.7, 1.08, 10.8, 0.35, subtitle, size=11.5, color=MUTED)


def add_footer(slide, text: str) -> None:
    add_text(slide, 10.9, 7.02, 1.8, 0.18, text, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullet_panel(slide, bullets: list[str], *, left: float, top: float, width: float, height: float, fill_color: str = WHITE) -> None:
    add_box(slide, left, top, width, height, fill_color=fill_color, line_color=BORDER)
    box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.18), Inches(width - 0.4), Inches(height - 0.28))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = rgb(INK)
        p.space_after = Pt(8)
        p.level = 0


def add_stat_card(slide, left: float, top: float, width: float, height: float, value: str, label: str, color: str = ACCENT) -> None:
    add_box(slide, left, top, width, height, fill_color=WHITE, line_color=BORDER)
    add_text(slide, left + 0.18, top + 0.12, width - 0.3, 0.34, value, size=24, bold=True, color=color)
    add_text(slide, left + 0.18, top + 0.55, width - 0.3, 0.26, label, size=11, color=MUTED)


def add_image(slide, image_name: str, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    path = ASSETS_DIR / image_name
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_callout(slide, left: float, top: float, width: float, height: float, heading: str, body: str, *, fill_color: str = ACCENT_SOFT, heading_color: str = ACCENT) -> None:
    add_box(slide, left, top, width, height, fill_color=fill_color, line_color=None)
    add_text(slide, left + 0.18, top + 0.12, width - 0.35, 0.28, heading, size=13, bold=True, color=heading_color)
    add_text(slide, left + 0.18, top + 0.42, width - 0.35, height - 0.5, body, size=11.5, color=INK)


def add_section_banner(slide, title: str, subtitle: str) -> None:
    add_box(slide, 0.7, 1.25, 12.0, 4.8, fill_color=ACCENT, line_color=None)
    add_text(slide, 1.1, 2.35, 8.2, 0.8, title, size=29, bold=True, color=WHITE)
    add_text(slide, 1.1, 3.25, 7.4, 0.7, subtitle, size=16, color="E6EBFF")
    add_box(slide, 9.25, 2.0, 2.55, 2.55, fill_color="FFFFFF", line_color=None)
    add_text(slide, 9.45, 2.58, 2.15, 1.2, "Project\nHighlights", size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def add_full_figure_slide(prs: Presentation, *, title: str, subtitle: str, image_name: str, takeaway_heading: str, takeaway_body: str, footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    add_box(slide, 0.72, 1.55, 8.9, 5.35, fill_color=WHITE, line_color=BORDER)
    add_image(slide, image_name, 0.95, 1.8, 8.45, None)
    add_callout(slide, 9.95, 1.9, 2.65, 1.85, takeaway_heading, takeaway_body, fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 9.95, 4.05, 2.65, 1.75, "Why it matters", "This figure is one of the main pieces of evidence for the final project conclusion.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, footer)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, "EEF2FF")
    add_box(slide, 0, 0, 4.6, 7.5, fill_color=ACCENT, line_color=None, radius=False)
    add_text(slide, 0.8, 1.05, 7.0, 1.2, "Predicting Plastic Degradation\nwith Machine Learning", size=28, bold=True, color=INK)
    add_text(slide, 0.82, 2.45, 6.5, 0.5, "From descriptor baselines to reliability-aware model selection", size=15, color=MUTED)
    add_box(slide, 0.82, 4.95, 3.2, 1.35, fill_color=WHITE, line_color=None)
    add_text(slide, 1.05, 5.28, 2.8, 0.28, "Aryamann Zutshi", size=16, bold=True, color=ACCENT)
    add_text(slide, 1.05, 5.66, 2.8, 0.24, "EECS 692", size=11.5, color=MUTED)
    add_text(slide, 1.05, 5.92, 2.8, 0.24, "Instructor: Dr. David Johnson", size=11.5, color=MUTED)
    add_text(slide, 8.15, 1.35, 3.7, 0.38, "Core idea", size=13, bold=True, color=ACCENT)
    add_text(slide, 8.15, 1.78, 4.2, 1.4, "This project moved from simply building classifiers to asking a more important question: which model is actually trustworthy when conditions get harder?", size=18, color=INK)
    add_callout(slide, 8.1, 4.15, 3.85, 1.45, "Final recommendation", "Top-ranked feature set + Random Forest gave the best balance between accuracy, calibration, and robustness.", fill_color=WHITE)
    add_footer(slide, "Slide 1")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Project Objective", "The project changed from performance-first modeling to reliability-aware selection.")
    add_bullet_panel(
        slide,
        [
            "Predict biodegradability from molecular descriptor data",
            "Compare classical ML, neural networks, and graph-inspired methods",
            "Stress-test the pipeline with validation, feature engineering, and uncertainty analysis",
            "Choose a final model that is accurate and trustworthy",
        ],
        left=0.75,
        top=1.55,
        width=6.0,
        height=4.9,
    )
    add_callout(slide, 7.15, 1.6, 5.25, 1.55, "Research shift", "The most important change in the project was moving from “which model scores highest?” to “which model can actually be defended?”")
    add_callout(slide, 7.15, 3.4, 5.25, 1.3, "Why this matters", "A high score alone does not show whether the model generalizes, calibrates well, or knows when it might be wrong.", fill_color="FFF3E8", heading_color=ORANGE)
    add_callout(slide, 7.15, 4.95, 5.25, 1.3, "Final framing", "The final model should be selected by reliability, not just by a single favorable metric.", fill_color="E7F8F4", heading_color=TEAL)
    add_footer(slide, "Slide 2")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Dataset and Scope", "QSAR biodegradation dataset used throughout the project")
    add_stat_card(slide, 0.85, 1.55, 2.2, 1.0, "1,055", "samples")
    add_stat_card(slide, 3.25, 1.55, 2.2, 1.0, "41", "molecular descriptors", color=TEAL)
    add_stat_card(slide, 5.65, 1.55, 2.2, 1.0, "2", "target classes", color=ORANGE)
    add_bullet_panel(
        slide,
        [
            "Target: RB vs NRB",
            "Examples of descriptors: SpMax_L, SpMax_B_m, SM6_B_m, nN, LOC",
            "Good for QSAR classification, but not a true polymer graph dataset",
            "No continuous degradation-rate constants or true pathway labels",
        ],
        left=0.85,
        top=3.0,
        width=6.1,
        height=3.1,
    )
    add_callout(slide, 7.35, 2.0, 4.7, 1.45, "Key limitation", "The data supported biodegradability classification well, but it did not directly support every original project goal, so the later weeks had to be adapted honestly.")
    add_callout(slide, 7.35, 3.8, 4.7, 1.55, "What that changed", "Regression on true rate constants, pathway labels, and chemistry-native polymer graphs became future-work directions instead of claims I could make from this dataset.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, "Slide 3")

    # Slide 4 section banner
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_section_banner(slide, "How the Project Evolved", "Each week answered a different modeling question instead of repeating the same experiment.")
    add_footer(slide, "Slide 4")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Model Progression: Weeks 5 to 7", "Starting with classical baselines before testing more complex models")
    add_bullet_panel(
        slide,
        [
            "The baseline stage established the first strong reference point",
            "Random Forest outperformed Logistic Regression on the standard split",
            "Neural and descriptor-graph models tested whether added complexity would help",
            "They did not beat the strongest classical baseline",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.55,
    )
    add_callout(slide, 7.1, 1.8, 5.2, 1.25, "Main takeaway", "Model complexity alone did not improve results. Representation quality mattered more.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 3.35, 5.2, 1.15, "Transition", "This set up the next question: do these models still hold up once the data distribution changes?", fill_color="FFF3E8", heading_color=ORANGE)
    add_callout(slide, 7.1, 4.85, 5.2, 1.15, "Figure on next slide", "The next slide shows the direct comparison between the baseline, the FNN, and the descriptor-graph prototype.")
    add_footer(slide, "Slide 5")

    add_full_figure_slide(
        prs,
        title="Figure: Model Progression Comparison",
        subtitle="Direct comparison of Weeks 5 to 7 models on the standard evaluation setup",
        image_name="model_progression_comparison.png",
        takeaway_heading="What it showed",
        takeaway_body="Random Forest stayed strongest, while the FNN and descriptor-graph prototype did not justify their extra complexity.",
        footer="Slide 6",
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Cross-Environment Validation", "Cross-environment validation tested whether learned patterns transferred under distribution shift")
    add_bullet_panel(
        slide,
        [
            "I created 3 proxy environments using descriptor-space clustering",
            "Then I ran leave-one-environment-out validation",
            "All models dropped hard under the shifted distribution",
            "This was where generalization clearly became the central problem",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.55,
    )
    add_callout(slide, 7.1, 1.8, 5.2, 1.25, "Most important result", "Random Forest still performed best, but even it dropped to about 0.421 mean held-out accuracy.", fill_color="FFF3E8", heading_color=ORANGE)
    add_callout(slide, 7.1, 3.35, 5.2, 1.25, "What changed", "From this point on, the project was no longer just about prediction quality. It was about generalization and reliability.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 4.9, 5.2, 1.1, "Figure on next slide", "The next slide shows the full cross-environment comparison clearly on its own.")
    add_footer(slide, "Slide 7")

    # Slide 8
    add_full_figure_slide(
        prs,
        title="Figure: Cross-Environment Performance",
        subtitle="Leave-one-environment-out validation across proxy environments",
        image_name="cross_environment_validation.png",
        takeaway_heading="What it showed",
        takeaway_body="Even the strongest models struggled under distribution shift, which made robustness the real bottleneck in the project.",
        footer="Slide 8",
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Validation Repair and Feature Engineering", "Weeks 9 and 10 focused on pipeline correctness and chemistry-aware expansion")
    add_bullet_panel(
        slide,
        [
            "I repaired the validation setup with 5-fold stratified cross-validation",
            "SMOTE was applied only on training folds",
            "Then I added chemistry-aware proxy features for polarity, heteroatoms, and topology",
            "The best feature-engineering setup improved RB recall without changing the core dataset",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.7,
    )
    add_callout(slide, 7.1, 1.8, 5.2, 1.3, "Best feature-engineering result", "Tier 1 + Tier 2 with FNN and SMOTE reached 0.8758 accuracy and 0.8399 RB recall.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 3.45, 5.2, 1.2, "Why this mattered", "It showed that chemistry-aware features could help, but they did not remove the broader robustness problem.", fill_color="FFF3E8", heading_color=ORANGE)
    add_callout(slide, 7.1, 4.95, 5.2, 1.0, "Figure on next slide", "The next slide isolates the feature-engineering comparison chart.")
    add_footer(slide, "Slide 9")

    # Slide 10
    add_full_figure_slide(
        prs,
        title="Figure: Feature Engineering",
        subtitle="Comparing baseline descriptors against the chemistry-aware proxy feature expansion",
        image_name="feature_engineering_comparison.png",
        takeaway_heading="What it showed",
        takeaway_body="The added chemistry-aware features helped in some settings, especially for recall, but they did not solve the deeper generalization issue.",
        footer="Slide 10",
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Feature Selection Changed the Story", "Feature selection showed that smaller, more meaningful feature sets could stay competitive")
    add_bullet_panel(
        slide,
        [
            "The strongest ranked features included SpMax_B_m, SpMax_L, SM6_B_m, and SpPosA_B_p",
            "Proxy features like mass_topology_proxy and polarity_proxy also mattered",
            "Top-ranked and reduced-hybrid stayed competitive",
            "Proxy-only features were clearly the weakest set",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.7,
    )
    add_callout(slide, 7.1, 1.9, 5.2, 1.3, "Why this mattered", "The top-ranked feature set did not just reduce dimensionality. It ended up supporting the strongest final reliability result in the final reliability analysis.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 3.55, 5.2, 1.15, "Figure on next slide", "The next slide isolates the feature-importance chart so the ranking is easy to read.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, "Slide 11")

    # Slide 12
    add_full_figure_slide(
        prs,
        title="Figure: Feature Importance",
        subtitle="Top-ranked features from Random Forest importance, permutation importance, and mutual information",
        image_name="feature_importance.png",
        takeaway_heading="What it showed",
        takeaway_body="A smaller, more focused feature set preserved strong signal and supported the strongest final recommendation.",
        footer="Slide 12",
    )

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Uncertainty and Calibration", "Uncertainty analysis asked whether the model knew when it might be wrong")
    add_bullet_panel(
        slide,
        [
            "I measured confidence, uncertainty, entropy, Brier score, ECE, and log loss",
            "Top-ranked was the best overall calibrated Random Forest feature set",
            "Wrong predictions were usually more uncertain in-distribution",
            "Under cross-environment shift, some models stayed confidently wrong",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.7,
    )
    add_callout(slide, 7.1, 1.85, 5.2, 1.25, "What changed", "At this point the project stopped being only about performance and became more about trustworthiness.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 3.4, 5.2, 1.25, "Main caution", "Distribution shift could still produce overconfident wrong predictions, which is exactly the kind of failure accuracy alone misses.", fill_color="FFF3E8", heading_color=ORANGE)
    add_callout(slide, 7.1, 4.95, 5.2, 1.0, "Figures on next slides", "The next two slides show the calibration and selective-prediction graphs separately.")
    add_footer(slide, "Slide 13")

    # Slide 14
    add_full_figure_slide(
        prs,
        title="Figure: Calibration Curve",
        subtitle="Comparing predicted confidence against observed biodegradability frequency",
        image_name="calibration_curve.png",
        takeaway_heading="What it showed",
        takeaway_body="The top-ranked feature set gave the best overall calibrated Random Forest behavior across the uncertainty analysis comparison.",
        footer="Slide 14",
    )

    # Slide 15
    add_full_figure_slide(
        prs,
        title="Figure: Selective Prediction",
        subtitle="Accuracy when keeping only the most confident predictions",
        image_name="selective_accuracy.png",
        takeaway_heading="What it showed",
        takeaway_body="Confidence could help improve trust, but it did not fully remove the cross-environment weakness.",
        footer="Slide 15",
    )

    # Slide 16
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Final Model Selection", "The final reliability analysis combined performance and trustworthiness into one recommendation")
    add_bullet_panel(
        slide,
        [
            "I ranked candidates using accuracy, calibration, uncertainty behavior, selective prediction, and cross-environment robustness",
            "The final choice was not made from one metric",
            "The recommendation came from the full reliability profile",
        ],
        left=0.75,
        top=1.55,
        width=5.9,
        height=4.3,
    )
    add_callout(slide, 7.1, 1.9, 5.15, 1.3, "Best overall candidate", "top_ranked | random_forest_classifier\nreliability score: 0.8860", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.1, 3.55, 5.15, 1.2, "Best raw cross-env accuracy", "top_ranked | logistic_regression at 0.6455", fill_color=WHITE, heading_color=ACCENT)
    add_callout(slide, 7.1, 5.05, 5.15, 1.05, "Figures on next slides", "The next two slides show the final scoreboard and the calibration-versus-accuracy tradeoff clearly on their own.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, "Slide 16")

    # Slide 17
    add_full_figure_slide(
        prs,
        title="Figure: Overall Reliability Scoreboard",
        subtitle="Composite reliability ranking across all model and feature-set candidates",
        image_name="model_reliability_scoreboard.png",
        takeaway_heading="What it showed",
        takeaway_body="Top-ranked Random Forest came out first once the project used reliability rather than only raw performance.",
        footer="Slide 17",
    )

    # Slide 18
    add_full_figure_slide(
        prs,
        title="Figure: Accuracy vs Calibration Tradeoff",
        subtitle="Comparing in-distribution and cross-environment tradeoffs across final candidates",
        image_name="accuracy_calibration_tradeoff.png",
        takeaway_heading="What it showed",
        takeaway_body="The strongest final model was the one that stayed competitive across several metrics instead of dominating only one.",
        footer="Slide 18",
    )

    # Slide 19
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Final Recommendation", "What I would defend as the final model choice from the whole project")
    add_box(slide, 0.9, 1.55, 5.0, 4.8, fill_color=ACCENT, line_color=None)
    add_text(slide, 1.25, 2.0, 4.2, 0.35, "Recommended final model", size=15, bold=True, color="DDE6FF")
    add_text(slide, 1.25, 2.45, 4.2, 1.2, "Top-ranked\nRandom Forest", size=28, bold=True, color=WHITE)
    add_text(slide, 1.25, 4.2, 4.15, 0.7, "This combination was not the winner in every single category, but it was the strongest overall once I looked at accuracy, calibration, uncertainty behavior, and shifted data performance together.", size=15, color="EAF0FF")
    add_bullet_panel(
        slide,
        [
            "Competitive standard accuracy",
            "Strong calibration for a high-performing model",
            "Useful uncertainty separation between right and wrong predictions",
            "Much better cross-environment behavior than the weaker feature sets",
        ],
        left=6.4,
        top=1.7,
        width=5.9,
        height=3.3,
    )
    add_callout(slide, 6.45, 5.25, 5.85, 1.1, "Bottom line", "The best final model was not just the highest-scoring model. It was the most defensible one.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, "Slide 19")

    # Slide 20
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Limitations and Next Steps", "What this project answered well, and what still needs stronger data")
    add_bullet_panel(
        slide,
        [
            "Descriptor-based dataset, not a true polymer graph dataset",
            "No real degradation-rate constants",
            "No pathway labels such as oxidation or hydrolysis",
            "Cross-environment robustness remains the biggest weakness",
        ],
        left=0.85,
        top=1.6,
        width=5.9,
        height=4.2,
    )
    add_callout(slide, 7.15, 1.85, 5.0, 1.45, "Best next step", "Move to a dataset with real polymer structure and stronger chemistry labels so that graph-based learning can be evaluated honestly.")
    add_callout(slide, 7.15, 3.65, 5.0, 1.45, "Validation next step", "Test the final model on a genuinely separate dataset or environment source instead of only proxy distribution shifts.", fill_color="E7F8F4", heading_color=TEAL)
    add_callout(slide, 7.15, 5.45, 5.0, 0.95, "Final lesson", "Trustworthiness ended up being the most important outcome of the project, not just raw model performance.", fill_color="FFF3E8", heading_color=ORANGE)
    add_footer(slide, "Slide 20")

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build()
    print(path)
